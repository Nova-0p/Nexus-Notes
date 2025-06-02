from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import tempfile
from datetime import datetime
from openai import OpenAI
import PyPDF2
from pptx import Presentation
import io
import logging
import markdown
from weasyprint import HTML, CSS
from markdown.extensions import codehilite, tables, toc
import re

from dotenv import load_dotenv
load_dotenv()  # This loads the variables from .env

# Initialize Flask app
app = Flask(__name__)
CORS(app)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Custom prompts for each summary type
PROMPTS = {
   "exam_focused": """
You are an expert academic assistant specializing in exam preparation. I need you to create a comprehensive exam study summary from the following document content. make sure the output is clean ore any wweird formatting and easily readable. Please make sure to include all the key points and main ideas from the document

Document Name: {filename}

Please create a detailed study summary that includes:

## 📚 EXAM STUDY GUIDE

### 1. **Key Topics & Concepts**
- List and explain the main topics covered
- Organize by importance and frequency of appearance in exams

### 2. **Important Definitions**
- Define key terms and concepts clearly
- Provide examples where applicable
- Create memory aids or mnemonics

### 3. **Critical Points for Exam Success**
- Highlight the most important information for exam preparation
- Mark high-priority items with ⭐
- Include common exam patterns

### 4. **Formulas & Equations** (if applicable)
- List all formulas clearly with explanations
- Provide example problems and solutions
- Include units and conditions

### 5. **Examples & Case Studies**
- Summarize important examples with step-by-step solutions
- Explain the logic behind each solution
- Provide similar practice scenarios

### 6. **Study Tips & Exam Strategy**
- Suggest what to focus on for maximum marks
- Time management recommendations
- Common mistakes to avoid

### 7. **Quick Review Section**
- Bullet-point summary for last-minute revision
- Key facts and figures
- Important dates, names, or events

### 8. **Practice Questions**
- Generate 5-10 potential exam questions based on content
- Provide brief answer outlines

Format the response in a clear, organized manner optimized for exam success.  Avoid symbols or garbled text."

Document Content: {content}
""",

    "research_summary": """
You are my AI research assistant. I will provide academic content such as papers, studies, or literature reviews. Please deliver a comprehensive research summary with the following structure:

## 🔍 RESEARCH ANALYSIS SUMMARY

### 1. **Research Overview**
- Title, Authors, Publication details
- Research question or hypothesis
- Study purpose and objectives

### 2. **Methodology**
- Research type (qualitative/quantitative/mixed)
- Data collection methods, sample size
- Analysis tools or statistical techniques used

### 3. **Key Findings**
- Summarize major results with statistical significance
- Highlight breakthrough discoveries
- Compare findings with previous studies (comparison analysis)

### 4. **Implications & Impact**
- Practical or theoretical significance
- Real-world applications
- Industry or academic impact

### 5. **Critical Analysis**
- Strengths and weaknesses of the study
- Validity and reliability assessment
- Potential biases or limitations

### 6. **Future Research Directions**
- Suggestions for follow-up studies
- Unresolved questions
- Emerging trends

### 7. **Glossary**
- Key academic or domain-specific terms with definitions
- Technical terminology explained

### 8. **Visual Summary Concept**
- Describe key relationships or processes
- Suggest flowcharts or conceptual maps

⛔ Do NOT include speculative interpretations or unrelated fields.

Document: {filename}
Content: {content}
""",

    "business_analysis": """
You are my AI business analyst. I will share reports, briefs, or strategy documents. Summarize the content as a structured business insight document:

## 📈 BUSINESS ANALYSIS REPORT

### 1. **Executive Summary**
- Main goals, insights, and key decisions
- Strategic recommendations in brief
- Critical success factors

### 2. **Strategic Implications**
- Impact on operations, competition, and growth
- Market positioning analysis
- Competitive advantages identified

### 3. **Financial Impact Analysis**
- Cost-benefit breakdown with specific numbers
- ROI estimates and projections
- Revenue impact and financial risks

### 4. **Stakeholder Analysis**
- Key stakeholders and their interests
- Impact assessment for each group
- Communication requirements

### 5. **Risk Assessment & Mitigation**
- Identify potential risks (High/Medium/Low)
- Propose specific mitigation strategies
- Contingency planning recommendations

### 6. **Action Plan & Implementation**
- Specific action steps with priorities
- Timeline view (short-term, mid-term, long-term)
- Resource requirements and budget implications

### 7. **Performance Metrics Dashboard**
- Key Performance Indicators (KPIs)
- Success metrics and benchmarks
- Monitoring and evaluation criteria

### 8. **Business Terminology**
- Industry-specific terms and acronyms
- Financial and strategic concepts explained

⛔ Do NOT include motivational commentary or unrelated market data.

Document: {filename}
Content: {content}
""",

    "legal_summary": """
You are my legal document assistant. I will provide legal contracts, case summaries, or regulations. Present a professional legal brief:

## ⚖️ LEGAL DOCUMENT ANALYSIS

### 1. **Case/Document Overview**
- Title, parties involved, jurisdiction
- Document type and legal context
- Date and relevant legal framework

### 2. **Legal Issues & Key Clauses**
- **PRIMARY LEGAL ISSUES:** (in bold)
- **KEY CLAUSES:** Identify and explain critical provisions
- **COMPLIANCE REQUIREMENTS:** Mandatory obligations
- **LIABILITY CONCERNS:** Risk areas and exposures

### 3. **Arguments & Legal Reasoning**
- Summary of main legal arguments
- Supporting precedents and statutes
- Counter-arguments and rebuttals

### 4. **Judgment/Decision/Outcome**
- Final ruling or contractual terms
- Legal reasoning behind decisions
- Precedential value

### 5. **Legal Implications & Consequences**
- **IMMEDIATE OBLIGATIONS:** What must be done now
- **LONG-TERM IMPLICATIONS:** Future legal effects
- **PRECEDENT IMPACT:** Effect on similar cases

### 6. **Compliance Checklist**
- Required actions with specific deadlines
- Documentation requirements
- Filing obligations and procedures

### 7. **Legal Terminology**
- Latin phrases and legal concepts
- Jurisdiction-specific terms
- Procedural terminology

### 8. **Comparative Analysis**
- Similar cases or legal precedents
- Distinguishing factors
- Trend analysis in related law

⛔ Do NOT include legal advice or personal interpretation.

Document: {filename}
Content: {content}
""",

    "meeting_notes": """
You are my AI assistant for professional meeting notes. I will provide transcripts or recordings. Format the summary clearly as follows:

## 📝 PROFESSIONAL MEETING SUMMARY

### 1. **Meeting Overview**
- **Date:** {current_date}
- **Participants:** (Extract from content)
- **Purpose:** Main objectives and agenda
- **Duration:** (if available)

### 2. **Executive Summary**
- **KEY OUTCOMES:** Major decisions and agreements
- **URGENT HIGHLIGHTS:** Time-sensitive items
- **STRATEGIC DECISIONS:** Important directional choices

### 3. **Topic-Based Discussion Breakdown**

#### **Topic 1: [Extract Topic Name]**
- **Discussion Points:**
  • Key arguments presented
  • Different perspectives shared
  • Concerns raised
- **Decision:** Final outcome or agreement
- **Context:** Background information

#### **Topic 2: [Extract Topic Name]**
[Repeat structure for each major topic]

### 4. **Action Items & Responsibilities**

| **Task** | **Owner** | **Deadline** | **Priority** | **Status** |
|----------|-----------|--------------|--------------|------------|
| [Extract from content] | [Name] | [Date] | High/Med/Low | Pending |

### 5. **Follow-Up Requirements**
- **CLARIFICATIONS NEEDED:** Open questions
- **APPROVALS REQUIRED:** Pending decisions
- **DEPENDENCIES:** Items waiting on external factors

### 6. **Next Meeting Agenda**
- **CARRY-OVER ITEMS:** Unresolved topics
- **NEW TOPICS:** Emerging issues
- **PREPARATION REQUIRED:** Pre-meeting tasks

### 7. **Key Terminology & Acronyms**
- Project-specific terms used
- Internal company language
- Technical acronyms explained

⛔ Do NOT include informal conversations or off-topic discussions.

Document: {filename}
Content: {content}
""",

    "technical_documentation": """
You are my technical documentation assistant. I will provide system designs, source code, or architecture notes. Return a developer-ready document with:

## 💻 TECHNICAL DOCUMENTATION

### 1. **Technical Overview**
- **Purpose:** System/component objective
- **Technology Stack:** Languages, frameworks, tools
- **Architecture Pattern:** Design approach used
- **Key Features:** Main functionalities

### 2. **System Architecture**
```
[Text-based Architecture Diagram]
┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   Frontend  │ ←→ │   Backend   │ ←→ │  Database   │
└─────────────┘    └─────────────┘    └─────────────┘
```
- **Data Flow:** How information moves through system
- **Component Interactions:** Key relationships
- **Security Layers:** Authentication and authorization

### 3. **Core Modules & Components**
#### **Module 1: [Name]**
- **Responsibility:** What it does
- **Dependencies:** Required components
- **Interfaces:** Input/output specifications

#### **Module 2: [Name]**
[Repeat for each major component]

### 4. **API Reference**
#### **Endpoint: [Method] /api/endpoint**
- **Parameters:**
  - `param1` (string, required): Description
  - `param2` (integer, optional): Description
- **Response:** Expected return format
- **Example:**
```json
{
  "status": "success",
  "data": {...}
}
```

### 5. **Setup & Installation**
```bash
# Step 1: Environment setup
# Step 2: Dependency installation
# Step 3: Configuration
# Step 4: Running the application
```

### 6. **Usage Examples**
```python
# Example 1: Basic usage
# Example 2: Advanced configuration
# Example 3: Error handling
```

### 7. **Troubleshooting & FAQs**
- **Common Issues:** Known problems and solutions
- **Performance Tips:** Optimization recommendations
- **Debugging Guide:** How to diagnose problems

### 8. **Technical Glossary**
- APIs, protocols, and architectural terms
- Framework-specific terminology
- Performance and security concepts

⛔ Do NOT include speculative performance metrics or outdated technologies.

Document: {filename}
Content: {content}
""",

    "medical_summary": """
You are my clinical AI assistant. I will provide clinical notes, case studies, or textbook content. Create a structured, study-focused summary:

## 🩺 MEDICAL CASE ANALYSIS

### 1. **Case Overview / Condition Summary**
- **Patient Profile:** Demographics and relevant history
- **Chief Complaint:** Primary symptoms presented
- **History of Present Illness:** Timeline and progression
- **Relevant Medical History:** Significant past conditions

### 2. **Clinical Assessment**
- **Physical Examination Findings:** Key signs and symptoms
- **Diagnostic Criteria:** How diagnosis was established
- **Laboratory Values:** Significant lab results with normal ranges
- **Differential Diagnosis:** Alternative conditions considered

### 3. **Treatment Plan & Management**
- **Immediate Interventions:** Emergency or urgent treatments
- **Medications:** Drugs, dosages, and administration routes
- **Procedures:** Surgical or therapeutic interventions
- **Monitoring:** What to watch for and when

### 4. **Pathophysiology Explanation**
- **Disease Mechanism:** How the condition develops
- **Physiological Impact:** Body system effects
- **Progression Pattern:** Expected disease course

### 5. **Medical Study Guide Features**
- **Memory Aids:** Mnemonics for symptoms, treatments
- **Common Errors:** Frequent diagnostic mistakes to avoid
- **Board Exam Points:** High-yield information for exams

### 6. **Clinical Decision Making**
```
Patient Presentation
        ↓
Initial Assessment
        ↓
Diagnostic Workup
        ↓
Treatment Selection
        ↓
Monitoring & Follow-up
```

### 7. **Medical Terminology**
- **Abbreviations:** Medical acronyms with full forms
- **Drug Names:** Generic and brand names
- **Anatomical Terms:** Relevant body structures
- **Pathological Terms:** Disease-specific vocabulary

### 8. **Quick Review Cards**
- **Symptoms:** Key presenting features
- **Treatment:** First-line therapies
- **Prognosis:** Expected outcomes
- **Complications:** Potential adverse effects

⛔ Do NOT include non-evidence-based remedies or personal health advice.

Document: {filename}
Content: {content}
"""
}

# CSS styles for PDF generation
PDF_CSS = """
@page {
    size: A4;
    margin: 2cm;
    @top-center {
        content: "Summary Document";
        font-size: 10pt;
        color: #666;
    }
    @bottom-center {
        content: counter(page);
        font-size: 10pt;
        color: #666;
    }
}

body {
    font-family: 'Arial', 'Helvetica', sans-serif;
    line-height: 1.6;
    color: #333;
    font-size: 11pt;
    margin: 0;
    padding: 0;
}

h1 {
    color: #2c3e50;
    font-size: 24pt;
    margin-bottom: 20px;
    padding-bottom: 10px;
    border-bottom: 3px solid #3498db;
    page-break-after: avoid;
}

h2 {
    color: #34495e;
    font-size: 18pt;
    margin-top: 25px;
    margin-bottom: 15px;
    page-break-after: avoid;
}

h3 {
    color: #5d6d7e;
    font-size: 14pt;
    margin-top: 20px;
    margin-bottom: 10px;
    page-break-after: avoid;
}

h4, h5, h6 {
    color: #85929e;
    font-size: 12pt;
    margin-top: 15px;
    margin-bottom: 8px;
    page-break-after: avoid;
}

p {
    margin-bottom: 12px;
    text-align: justify;
    orphans: 3;
    widows: 3;
}

ul, ol {
    margin-bottom: 15px;
    padding-left: 25px;
}

li {
    margin-bottom: 5px;
    page-break-inside: avoid;
}

strong {
    color: #2c3e50;
    font-weight: bold;
}

em {
    font-style: italic;
    color: #5d6d7e;
}

blockquote {
    margin: 20px 0;
    padding: 15px 20px;
    background-color: #f8f9fa;
    border-left: 5px solid #3498db;
    font-style: italic;
    page-break-inside: avoid;
}

code {
    background-color: #f1f2f6;
    padding: 2px 4px;
    border-radius: 3px;
    font-family: 'Courier New', monospace;
    font-size: 10pt;
}

pre {
    background-color: #f8f9fa;
    padding: 15px;
    border-radius: 5px;
    overflow-x: auto;
    margin: 15px 0;
    page-break-inside: avoid;
}

table {
    width: 100%;
    border-collapse: collapse;
    margin: 20px 0;
    page-break-inside: avoid;
}

th, td {
    border: 1px solid #ddd;
    padding: 8px;
    text-align: left;
}

th {
    background-color: #f2f2f2;
    font-weight: bold;
    color: #2c3e50;
}

hr {
    margin: 25px 0;
    border: none;
    border-top: 2px solid #ecf0f1;
    page-break-after: avoid;
}

.page-break {
    page-break-before: always;
}

.no-break {
    page-break-inside: avoid;
}
"""

def extract_text_from_pdf(file_stream):
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting PDF text: {str(e)}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_pptx(file_stream):
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {str(e)}")
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")

def extract_text_from_txt(file_stream):
    """Extract text from text file"""
    try:
        content = file_stream.read()
        if isinstance(content, bytes):
            content = content.decode('utf-8')
        return content.strip()
    except Exception as e:
        logger.error(f"Error extracting TXT text: {str(e)}")
        raise Exception(f"Failed to extract text from text file: {str(e)}")

def clean_markdown_text(text):
    """Clean and prepare markdown text for better PDF conversion"""
    # Remove excessive newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Fix heading spacing
    text = re.sub(r'\n(#{1,6})', r'\n\n\1', text)
    text = re.sub(r'(#{1,6}[^\n]*)\n([^\n#])', r'\1\n\n\2', text)
    
    # Ensure proper list formatting
    text = re.sub(r'\n([*+-])', r'\n\n\1', text)
    text = re.sub(r'\n(\d+\.)', r'\n\n\1', text)
    
    # Clean up bold and italic formatting
    text = re.sub(r'\*{3,}([^*]+)\*{3,}', r'**\1**', text)
    text = re.sub(r'\*{2,}([^*]+)\*{2,}', r'**\1**', text)
    
    return text.strip()

def markdown_to_pdf(markdown_text, output_path):
    """Convert markdown text to PDF"""
    try:
        # Clean the markdown text
        cleaned_markdown = clean_markdown_text(markdown_text)
        
        # Convert markdown to HTML
        md = markdown.Markdown(
            extensions=[
                'markdown.extensions.tables',
                'markdown.extensions.fenced_code',
                'markdown.extensions.codehilite',
                'markdown.extensions.toc',
                'markdown.extensions.nl2br'
            ],
            extension_configs={
                'codehilite': {
                    'css_class': 'highlight'
                }
            }
        )
        
        html_content = md.convert(cleaned_markdown)
        
        # Create complete HTML document
        full_html = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Summary Document</title>
</head>
<body>
    {html_content}
</body>
</html>
        """
        
        # Convert HTML to PDF using WeasyPrint
        html_doc = HTML(string=full_html)
        css_doc = CSS(string=PDF_CSS)
        
        html_doc.write_pdf(output_path, stylesheets=[css_doc])
        logger.info(f"PDF generated successfully: {output_path}")
        
    except Exception as e:
        logger.error(f"Error converting markdown to PDF: {str(e)}")
        raise Exception(f"Failed to generate PDF: {str(e)}")

def generate_summary_with_openai(content, filename, prompt_type):
    """Generate summary using OpenAI API"""
    try:
        # Get the appropriate prompt template
        if prompt_type not in PROMPTS:
            raise ValueError(f"Invalid prompt type: {prompt_type}")
        
        prompt_template = PROMPTS[prompt_type]
        
        # Format the prompt with content and filename
        current_date = datetime.now().strftime("%Y-%m-%d")
        formatted_prompt = prompt_template.format(
            content=content,
            filename=filename,
            current_date=current_date
        )
        
        # Make API call to OpenAI
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # You can change this to gpt-4 if you have access
            messages=[
                {
                    "role": "system", 
                    "content": "You are an expert AI assistant that creates detailed, structured summaries using markdown format. Use proper markdown syntax with headers (#, ##, ###), bold (**text**), italic (*text*), bullet points (-), and other markdown elements to create well-formatted, readable summaries."
                },
                {
                    "role": "user", 
                    "content": formatted_prompt
                }
            ],
            max_tokens=4000,
            temperature=0.3,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0
        )
        
        summary = response.choices[0].message.content
        return summary
        
    except Exception as e:
        logger.error(f"Error generating summary with OpenAI: {str(e)}")
        raise Exception(f"Failed to generate summary: {str(e)}")

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Get prompt type and output format
        prompt_type = request.form.get('prompt_type', 'exam_focused')
        output_format = request.form.get('output_format', 'pdf')  # 'pdf' or 'text'
        
        if prompt_type not in PROMPTS:
            return jsonify({'error': 'Invalid prompt type'}), 400
        
        # Check file size (limit to 10MB)
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > 10 * 1024 * 1024:  # 10MB limit
            return jsonify({'error': 'File size too large. Maximum 10MB allowed.'}), 400
        
        # Extract text based on file type
        filename = file.filename.lower()
        file_stream = io.BytesIO(file.read())
        
        try:
            if filename.endswith('.pdf'):
                text_content = extract_text_from_pdf(file_stream)
            elif filename.endswith(('.ppt', '.pptx')):
                text_content = extract_text_from_pptx(file_stream)
            elif filename.endswith('.txt'):
                file_stream.seek(0)
                text_content = extract_text_from_txt(file_stream)
            else:
                return jsonify({'error': 'Unsupported file type. Please upload PDF, PowerPoint, or text files.'}), 400
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            return jsonify({'error': f'Failed to extract text from file: {str(e)}'}), 400
        
        # Check if text content is not empty
        if not text_content or len(text_content.strip()) < 10:
            return jsonify({'error': 'File appears to be empty or contains insufficient text content.'}), 400
        
        # Truncate content if too long (to stay within token limits)
        max_content_length = 15000  # Adjust based on your needs
        if len(text_content) > max_content_length:
            text_content = text_content[:max_content_length] + "\n\n[Content truncated due to length...]"
        
        # Generate summary using OpenAI
        try:
            summary = generate_summary_with_openai(text_content, file.filename, prompt_type)
        except Exception as e:
            logger.error(f"Error generating summary: {str(e)}")
            return jsonify({'error': f'Failed to generate summary: {str(e)}'}), 500
        
        # Return based on requested format
        if output_format == 'pdf':
            try:
                # Create temporary file for PDF
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                    pdf_path = temp_pdf.name
                
                # Convert markdown to PDF
                markdown_to_pdf(summary, pdf_path)
                
                # Generate filename for download
                original_name = os.path.splitext(file.filename)[0]
                pdf_filename = f"{original_name}_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                
                def remove_file(response):
                    try:
                        os.unlink(pdf_path)
                    except Exception:
                        pass
                    return response
                
                # Send PDF file
                return send_file(
                    pdf_path,
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name=pdf_filename
                )
                
            except Exception as e:
                logger.error(f"Error generating PDF: {str(e)}")
                return jsonify({'error': f'Failed to generate PDF: {str(e)}'}), 500
        else:
            # Return as plain text (original functionality)
            return summary, 200, {'Content-Type': 'text/plain; charset=utf-8'}
        
    except Exception as e:
        logger.error(f"Unexpected error in upload_file: {str(e)}")
        return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'message': 'Nexus Notes API is running',
        'timestamp': datetime.now().isoformat()
    })

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 10MB.'}), 413

@app.errorhandler(500)
def internal_server_error(e):
    logger.error(f"Internal server error: {str(e)}")
    return jsonify({'error': 'Internal server error. Please try again later.'}), 500

if __name__ == '__main__':
    # Check if OpenAI API key is set
    if not os.getenv('OPENAI_API_KEY'):
        print("ERROR: OPENAI_API_KEY environment variable is not set!")
        print("Please set your OpenAI API key as an environment variable:")
        print("export OPENAI_API_KEY='your-api-key-here'")
        exit(1)
    
    print("Starting Nexus Notes API server...")
    print("Make sure to set your OPENAI_API_KEY environment variable")
    print("PDF generation enabled - install required packages:")
    print("pip install markdown weasyprint")
    app.run(host='0.0.0.0', port=5000, debug=True)