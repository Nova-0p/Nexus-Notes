from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import tempfile
from datetime import datetime
from openai import OpenAI
import PyPDF2
from pptx import Presentation
import io
import logging
import markdown
from weasyprint import HTML, CSS
from weasyprint.text.fonts import FontConfiguration
import gc
import atexit
import threading
import time

from dotenv import load_dotenv
load_dotenv()  # This loads the variables from .env

# Initialize Flask app
app = Flask(__name__)
CORS(app)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Track temporary files for cleanup
temp_files = set()
temp_files_lock = threading.Lock()

def cleanup_temp_files():
    """Clean up all temporary files"""
    with temp_files_lock:
        for file_path in list(temp_files):
            try:
                if os.path.exists(file_path):
                    os.unlink(file_path)
                    logger.info(f"Cleaned up temp file: {file_path}")
                temp_files.discard(file_path)
            except Exception as e:
                logger.warning(f"Could not delete temp file {file_path}: {str(e)}")

# Register cleanup function to run on exit
atexit.register(cleanup_temp_files)

def add_temp_file(file_path):
    """Add a temporary file to the cleanup list"""
    with temp_files_lock:
        temp_files.add(file_path)

def remove_temp_file(file_path):
    """Remove and clean up a temporary file"""
    try:
        if os.path.exists(file_path):
            os.unlink(file_path)
        with temp_files_lock:
            temp_files.discard(file_path)
    except Exception as e:
        logger.warning(f"Could not delete temp file {file_path}: {str(e)}")

# Custom prompts for each summary type
PROMPTS = {
    "exam_focused": """
You are an expert academic assistant specializing in exam preparation. I need you to create a comprehensive exam study summary from the following document content. make sure the output is clean ore any wweird formatting and easily readable. Please make sure to include all the key points and main ideas from the document dont miss any topic and explain all topics completly with the same promppt below 

Document Name: {filename}

Please create a detailed study summary that includes:

## 📚 EXAM STUDY GUIDE

### 1. **Key Topics & Concepts**
- List and explain the main topics covered
- Organize by importance and frequency of appearance in exams

### 2. **Important Definitions**
- Define key terms and concepts clearly
- Provide examples where applicable
- Create memory aids or mnemonics

### 3. **Critical Points for Exam Success**
- Highlight the most important information for exam preparation
- Mark high-priority items with ⭐
- Include common exam patterns

### 4. **Formulas & Equations** (if applicable)
- List all formulas clearly with explanations
- Provide example problems and solutions
- Include units and conditions

### 5. **Examples & Case Studies**
- Summarize important examples with step-by-step solutions
- Explain the logic behind each solution
- Provide similar practice scenarios

### 6. **Study Tips & Exam Strategy**
- Suggest what to focus on for maximum marks
- Time management recommendations
- Common mistakes to avoid

### 7. **Quick Review Section**
- Bullet-point summary for last-minute revision
- Key facts and figures
- Important dates, names, or events

### 8. **Practice Questions**
- Generate 5-10 potential exam questions based on content
- Provide brief answer outlines

Format the response in a clear, organized manner optimized for exam success.  Avoid symbols or garbled text."

Document Content: {content}
""",

    "research_summary": """
You are my AI research assistant. I will provide academic content such as papers, studies, or literature reviews. Please deliver a comprehensive research summary with the following structure:dont miss any topic and explain all topics completly with the same promppt below 

## 🔍 RESEARCH ANALYSIS SUMMARY

### 1. **Research Overview**
- Title, Authors, Publication details
- Research question or hypothesis
- Study purpose and objectives

### 2. **Methodology**
- Research type (qualitative/quantitative/mixed)
- Data collection methods, sample size
- Analysis tools or statistical techniques used

### 3. **Key Findings**
- Summarize major results with statistical significance
- Highlight breakthrough discoveries
- Compare findings with previous studies (comparison analysis)

### 4. **Implications & Impact**
- Practical or theoretical significance
- Real-world applications
- Industry or academic impact

### 5. **Critical Analysis**
- Strengths and weaknesses of the study
- Validity and reliability assessment
- Potential biases or limitations

### 6. **Future Research Directions**
- Suggestions for follow-up studies
- Unresolved questions
- Emerging trends

### 7. **Glossary**
- Key academic or domain-specific terms with definitions
- Technical terminology explained

### 8. **Visual Summary Concept**
- Describe key relationships or processes
- Suggest flowcharts or conceptual maps

⛔ Do NOT include speculative interpretations or unrelated fields.

Document: {filename}
Content: {content}
""",

    "business_analysis": """
You are my AI business analyst. I will share reports, briefs, or strategy documents. Summarize the content as a structured business insight document:dont miss any topic and explain all topics completly with the same promppt below dont miss any topic and explain all topics completly with the same promppt below 

## 📈 BUSINESS ANALYSIS REPORT

### 1. **Executive Summary**
- Main goals, insights, and key decisions
- Strategic recommendations in brief
- Critical success factors

### 2. **Strategic Implications**
- Impact on operations, competition, and growth
- Market positioning analysis
- Competitive advantages identified

### 3. **Financial Impact Analysis**
- Cost-benefit breakdown with specific numbers
- ROI estimates and projections
- Revenue impact and financial risks

### 4. **Stakeholder Analysis**
- Key stakeholders and their interests
- Impact assessment for each group
- Communication requirements

### 5. **Risk Assessment & Mitigation**
- Identify potential risks (High/Medium/Low)
- Propose specific mitigation strategies
- Contingency planning recommendations

### 6. **Action Plan & Implementation**
- Specific action steps with priorities
- Timeline view (short-term, mid-term, long-term)
- Resource requirements and budget implications

### 7. **Performance Metrics Dashboard**
- Key Performance Indicators (KPIs)
- Success metrics and benchmarks
- Monitoring and evaluation criteria

### 8. **Business Terminology**
- Industry-specific terms and acronyms
- Financial and strategic concepts explained

⛔ Do NOT include motivational commentary or unrelated market data.

Document: {filename}
Content: {content}
""",

    "legal_summary": """
You are my legal document assistant. I will provide legal contracts, case summaries, or regulations. Present a professional legal brief:
dont miss any topic and explain all topics completly with the same promppt below 
## ⚖️ LEGAL DOCUMENT ANALYSIS

### 1. **Case/Document Overview**
- Title, parties involved, jurisdiction
- Document type and legal context
- Date and relevant legal framework

### 2. **Legal Issues & Key Clauses**
- **PRIMARY LEGAL ISSUES:** (in bold)
- **KEY CLAUSES:** Identify and explain critical provisions
- **COMPLIANCE REQUIREMENTS:** Mandatory obligations
- **LIABILITY CONCERNS:** Risk areas and exposures

### 3. **Arguments & Legal Reasoning**
- Summary of main legal arguments
- Supporting precedents and statutes
- Counter-arguments and rebuttals

### 4. **Judgment/Decision/Outcome**
- Final ruling or contractual terms
- Legal reasoning behind decisions
- Precedential value

### 5. **Legal Implications & Consequences**
- **IMMEDIATE OBLIGATIONS:** What must be done now
- **LONG-TERM IMPLICATIONS:** Future legal effects
- **PRECEDENT IMPACT:** Effect on similar cases

### 6. **Compliance Checklist**
- Required actions with specific deadlines
- Documentation requirements
- Filing obligations and procedures

### 7. **Legal Terminology**
- Latin phrases and legal concepts
- Jurisdiction-specific terms
- Procedural terminology

### 8. **Comparative Analysis**
- Similar cases or legal precedents
- Distinguishing factors
- Trend analysis in related law

⛔ Do NOT include legal advice or personal interpretation.

Document: {filename}
Content: {content}
""",

    "meeting_notes": """
You are my AI assistant for professional meeting notes. I will provide transcripts or recordings. Format the summary clearly as follows:
dont miss any topic and explain all topics completly with the same promppt below 
## 📝 PROFESSIONAL MEETING SUMMARY

### 1. **Meeting Overview**
- **Date:** {current_date}
- **Participants:** (Extract from content)
- **Purpose:** Main objectives and agenda
- **Duration:** (if available)

### 2. **Executive Summary**
- **KEY OUTCOMES:** Major decisions and agreements
- **URGENT HIGHLIGHTS:** Time-sensitive items
- **STRATEGIC DECISIONS:** Important directional choices

### 3. **Topic-Based Discussion Breakdown**

#### **Topic 1: [Extract Topic Name]**
- **Discussion Points:**
  • Key arguments presented
  • Different perspectives shared
  • Concerns raised
- **Decision:** Final outcome or agreement
- **Context:** Background information

#### **Topic 2: [Extract Topic Name]**
[Repeat structure for each major topic]

### 4. **Action Items & Responsibilities**

| **Task** | **Owner** | **Deadline** | **Priority** | **Status** |
|----------|-----------|--------------|--------------|------------|
| [Extract from content] | [Name] | [Date] | High/Med/Low | Pending |

### 5. **Follow-Up Requirements**
- **CLARIFICATIONS NEEDED:** Open questions
- **APPROVALS REQUIRED:** Pending decisions
- **DEPENDENCIES:** Items waiting on external factors

### 6. **Next Meeting Agenda**
- **CARRY-OVER ITEMS:** Unresolved topics
- **NEW TOPICS:** Emerging issues
- **PREPARATION REQUIRED:** Pre-meeting tasks

### 7. **Key Terminology & Acronyms**
- Project-specific terms used
- Internal company language
- Technical acronyms explained

⛔ Do NOT include informal conversations or off-topic discussions.

Document: {filename}
Content: {content}
""",

    "technical_documentation": """
You are my technical documentation assistant. I will provide system designs, source code, or architecture notes. Return a developer-ready document with:

## 💻 TECHNICAL DOCUMENTATION
dont miss any topic and explain all topics completly with the same promppt below 
### 1. **Technical Overview**
- **Purpose:** System/component objective
- **Technology Stack:** Languages, frameworks, tools
- **Architecture Pattern:** Design approach used
- **Key Features:** Main functionalities

### 2. **System Architecture**
```
[Text-based Architecture Diagram]
┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│   Frontend  │ ←→ │   Backend   │ ←→ │  Database   │
└─────────────┘    └─────────────┘    └─────────────┘
```
- **Data Flow:** How information moves through system
- **Component Interactions:** Key relationships
- **Security Layers:** Authentication and authorization

### 3. **Core Modules & Components**
#### **Module 1: [Name]**
- **Responsibility:** What it does
- **Dependencies:** Required components
- **Interfaces:** Input/output specifications

#### **Module 2: [Name]**
[Repeat for each major component]

### 4. **API Reference**
#### **Endpoint: [Method] /api/endpoint**
- **Parameters:**
  - `param1` (string, required): Description
  - `param2` (integer, optional): Description
- **Response:** Expected return format
- **Example:**
```json
{
  "status": "success",
  "data": {...}
}
```

### 5. **Setup & Installation**
```bash
# Step 1: Environment setup
# Step 2: Dependency installation
# Step 3: Configuration
# Step 4: Running the application
```

### 6. **Usage Examples**
```python
# Example 1: Basic usage
# Example 2: Advanced configuration
# Example 3: Error handling
```

### 7. **Troubleshooting & FAQs**
- **Common Issues:** Known problems and solutions
- **Performance Tips:** Optimization recommendations
- **Debugging Guide:** How to diagnose problems

### 8. **Technical Glossary**
- APIs, protocols, and architectural terms
- Framework-specific terminology
- Performance and security concepts

⛔ Do NOT include speculative performance metrics or outdated technologies.

Document: {filename}
Content: {content}
""",

    "medical_summary": """
You are my clinical AI assistant. I will provide clinical notes, case studies, or textbook content. Create a structured, study-focused summary:

## 🩺 MEDICAL CASE ANALYSIS
dont miss any topic and explain all topics completly with the same promppt below 
### 1. **Case Overview / Condition Summary**
- **Patient Profile:** Demographics and relevant history
- **Chief Complaint:** Primary symptoms presented
- **History of Present Illness:** Timeline and progression
- **Relevant Medical History:** Significant past conditions

### 2. **Clinical Assessment**
- **Physical Examination Findings:** Key signs and symptoms
- **Diagnostic Criteria:** How diagnosis was established
- **Laboratory Values:** Significant lab results with normal ranges
- **Differential Diagnosis:** Alternative conditions considered

### 3. **Treatment Plan & Management**
- **Immediate Interventions:** Emergency or urgent treatments
- **Medications:** Drugs, dosages, and administration routes
- **Procedures:** Surgical or therapeutic interventions
- **Monitoring:** What to watch for and when

### 4. **Pathophysiology Explanation**
- **Disease Mechanism:** How the condition develops
- **Physiological Impact:** Body system effects
- **Progression Pattern:** Expected disease course

### 5. **Medical Study Guide Features**
- **Memory Aids:** Mnemonics for symptoms, treatments
- **Common Errors:** Frequent diagnostic mistakes to avoid
- **Board Exam Points:** High-yield information for exams

### 6. **Clinical Decision Making**
```
Patient Presentation
        ↓
Initial Assessment
        ↓
Diagnostic Workup
        ↓
Treatment Selection
        ↓
Monitoring & Follow-up
```

### 7. **Medical Terminology**
- **Abbreviations:** Medical acronyms with full forms
- **Drug Names:** Generic and brand names
- **Anatomical Terms:** Relevant body structures
- **Pathological Terms:** Disease-specific vocabulary

### 8. **Quick Review Cards**
- **Symptoms:** Key presenting features
- **Treatment:** First-line therapies
- **Prognosis:** Expected outcomes
- **Complications:** Potential adverse effects

⛔ Do NOT include non-evidence-based remedies or personal health advice.

Document: {filename}
Content: {content}
"""
}

def extract_text_from_pdf(file_stream):
    """Extract text from PDF file with proper resource management"""
    pdf_reader = None
    try:
        file_stream.seek(0)
        pdf_reader = PyPDF2.PdfReader(file_stream)
        text = ""
        
        # Limit pages to prevent memory issues
        max_pages = min(len(pdf_reader.pages), 100)  # Limit to 100 pages
        
        for i in range(max_pages):
            try:
                page = pdf_reader.pages[i]
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            except Exception as e:
                logger.warning(f"Could not extract text from page {i+1}: {str(e)}")
                continue
        
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting PDF text: {str(e)}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")
    finally:
        # Force garbage collection
        gc.collect()

def extract_text_from_pptx(file_stream):
    """Extract text from PowerPoint file with proper resource management"""
    prs = None
    try:
        file_stream.seek(0)
        prs = Presentation(file_stream)
        text = ""
        
        # Limit slides to prevent memory issues
        max_slides = min(len(prs.slides), 50)  # Limit to 50 slides
        
        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            except Exception as e:
                logger.warning(f"Could not extract text from slide {i+1}: {str(e)}")
                continue
        
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {str(e)}")
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")
    finally:
        # Force garbage collection
        gc.collect()

def extract_text_from_txt(file_stream):
    """Extract text from text file with proper encoding handling"""
    try:
        file_stream.seek(0)
        content = file_stream.read()
        
        if isinstance(content, bytes):
            # Try different encodings
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    content = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                # If all encodings fail, use utf-8 with error handling
                content = content.decode('utf-8', errors='replace')
        
        return content.strip()
    except Exception as e:
        logger.error(f"Error extracting TXT text: {str(e)}")
        raise Exception(f"Failed to extract text from text file: {str(e)}")

def safe_markdown_to_pdf(markdown_content, filename="summary"):
    """Convert markdown content to PDF with better error handling"""
    temp_pdf_path = None
    try:
        # Convert markdown to HTML
        html_content = markdown.markdown(
            markdown_content,
            extensions=['tables', 'fenced_code', 'toc', 'nl2br']
        )
        
        # Create a complete HTML document with simplified styling
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>{filename} Summary</title>
            <style>
                @page {{
                    margin: 2cm;
                    size: A4;
                }}
                body {{
                    font-family: Arial, sans-serif;
                    line-height: 1.6;
                    color: #333;
                    font-size: 12pt;
                }}
                h1 {{ color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px; }}
                h2 {{ color: #34495e; border-bottom: 1px solid #ecf0f1; padding-bottom: 5px; }}
                h3 {{ color: #34495e; }}
                p {{ margin-bottom: 12px; }}
                ul, ol {{ margin-bottom: 15px; padding-left: 25px; }}
                li {{ margin-bottom: 5px; }}
                strong {{ font-weight: bold; }}
                code {{ background-color: #f4f4f4; padding: 2px 4px; }}
                pre {{ background-color: #f4f4f4; padding: 15px; }}
                table {{ border-collapse: collapse; width: 100%; margin: 15px 0; }}
                th, td {{ border: 1px solid #ddd; padding: 8px; }}
                th {{ background-color: #f2f2f2; font-weight: bold; }}
            </style>
        </head>
        <body>
            {html_content}
        </body>
        </html>
        """
        
        # Create a temporary file for the PDF
        temp_pdf_fd, temp_pdf_path = tempfile.mkstemp(suffix='.pdf')
        os.close(temp_pdf_fd)  # Close the file descriptor
        
        add_temp_file(temp_pdf_path)
        
        # Convert HTML to PDF using WeasyPrint with minimal configuration
        try:
            HTML(string=full_html).write_pdf(temp_pdf_path)
        except Exception as weasy_error:
            logger.error(f"WeasyPrint error: {str(weasy_error)}")
            # Try with even simpler HTML
            simple_html = f"""
            <!DOCTYPE html>
            <html>
            <head><meta charset="UTF-8"><title>Summary</title></head>
            <body style="font-family: Arial; margin: 20px;">
                <pre style="white-space: pre-wrap; font-family: Arial;">{markdown_content}</pre>
            </body>
            </html>
            """
            HTML(string=simple_html).write_pdf(temp_pdf_path)
        
        return temp_pdf_path
            
    except Exception as e:
        if temp_pdf_path:
            remove_temp_file(temp_pdf_path)
        logger.error(f"Error converting markdown to PDF: {str(e)}")
        raise Exception(f"Failed to convert markdown to PDF: {str(e)}")

def generate_summary_with_openai(content, filename, prompt_type):
    """Generate summary using OpenAI API with improved chunking"""
    try:
        if prompt_type not in PROMPTS:
            raise ValueError(f"Invalid prompt type: {prompt_type}")
        
        # Limit content length to prevent excessive API usage
        max_content_length = 50000  # Reduced from 100000
        if len(content) > max_content_length:
            content = content[:max_content_length]
            logger.info(f"Content truncated to {max_content_length} characters")
        
        # Improved chunking logic
        max_chunk_size = 15000  # Increased chunk size, reduced overlap
        chunks = []
        
        if len(content) <= max_chunk_size:
            # If content is small enough, process in one go
            chunks = [content]
        else:
            # Split into chunks with minimal overlap
            overlap = 200  # Reduced overlap
            start = 0
            
            while start < len(content):
                end = min(start + max_chunk_size, len(content))
                
                # Try to break at sentence boundaries
                if end < len(content):
                    # Look for sentence end within last 500 characters
                    search_start = max(end - 500, start)
                    sentence_end = content.rfind('.', search_start, end)
                    if sentence_end > start:
                        end = sentence_end + 1
                
                chunk = content[start:end]
                chunks.append(chunk)
                
                if end >= len(content):
                    break
                    
                start = end - overlap
                
                # Prevent infinite loop
                if start <= 0 or len(chunks) > 10:  # Limit to 10 chunks max
                    break
        
        logger.info(f"Processing {len(chunks)} chunks for {filename}")
        
        # Process chunks
        summaries = []
        current_date = datetime.now().strftime("%Y-%m-%d")
        
        for i, chunk in enumerate(chunks):
            try:
                prompt_template = PROMPTS[prompt_type]
                
                response = client.chat.completions.create(
                    model="gpt-4o-mini",  # Use more cost-effective model
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant that creates concise, well-structured summaries."},
                        {"role": "user", "content": prompt_template.format(
                            content=chunk,
                            filename=f"{filename}" + (f" (Part {i+1}/{len(chunks)})" if len(chunks) > 1 else ""),
                            current_date=current_date
                        )}
                    ],
                    max_tokens=2000,  # Reduced token limit
                    temperature=0.3,
                    timeout=30  # Add timeout
                )
                
                summaries.append(response.choices[0].message.content)
                
                # Small delay to prevent rate limiting
                if len(chunks) > 1:
                    time.sleep(0.5)
                    
            except Exception as e:
                logger.error(f"Error processing chunk {i+1}: {str(e)}")
                summaries.append(f"[Error processing part {i+1} of the document]")
        
        # Combine summaries
        if len(summaries) == 1:
            return summaries[0]
        else:
            combined_summary = "\n\n---\n\n".join(summaries)
            
            # If we have multiple summaries, create a final consolidated version
            try:
                final_response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "Combine these partial summaries into one coherent, well-structured document. Remove redundancy and create a flowing narrative."},
                        {"role": "user", "content": f"Combine these {len(summaries)} summaries of {filename}:\n\n{combined_summary}"}
                    ],
                    max_tokens=3000,
                    temperature=0.3,
                    timeout=30
                )
                return final_response.choices[0].message.content
            except Exception as e:
                logger.error(f"Error creating final summary: {str(e)}")
                return combined_summary  # Return combined version if final summary fails
        
    except Exception as e:
        logger.error(f"Error generating summary: {str(e)}")
        raise

@app.route('/upload', methods=['POST'])
def upload_file():
    file_stream = None
    try:
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Get prompt type and output format
        prompt_type = request.form.get('prompt_type', 'exam_focused')
        output_format = request.form.get('output_format', 'text')
        
        if prompt_type not in PROMPTS:
            return jsonify({'error': 'Invalid prompt type'}), 400
        
        # Check file size (limit to 5MB to prevent system overload)
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > 5 * 1024 * 1024:  # 5MB limit
            return jsonify({'error': 'File size too large. Maximum 5MB allowed.'}), 400
        
        # Extract text based on file type
        filename = file.filename.lower()
        file_stream = io.BytesIO(file.read())
        
        try:
            if filename.endswith('.pdf'):
                text_content = extract_text_from_pdf(file_stream)
            elif filename.endswith(('.ppt', '.pptx')):
                text_content = extract_text_from_pptx(file_stream)
            elif filename.endswith('.txt'):
                text_content = extract_text_from_txt(file_stream)
            else:
                return jsonify({'error': 'Unsupported file type. Please upload PDF, PowerPoint, or text files.'}), 400
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            return jsonify({'error': f'Failed to extract text from file: {str(e)}'}), 400
        finally:
            if file_stream:
                file_stream.close()
        
        # Check if text content is not empty
        if not text_content or len(text_content.strip()) < 10:
            return jsonify({'error': 'File appears to be empty or contains insufficient text content.'}), 400
        
        # Generate summary using OpenAI
        try:
            summary = generate_summary_with_openai(text_content, file.filename, prompt_type)
        except Exception as e:
            logger.error(f"Error generating summary: {str(e)}")
            return jsonify({'error': f'Failed to generate summary: {str(e)}'}), 500
        
        # Return based on requested format
        if output_format == 'pdf':
            pdf_path = None
            try:
                # Convert markdown summary to PDF
                base_filename = os.path.splitext(file.filename)[0]
                pdf_path = safe_markdown_to_pdf(summary, base_filename)
                
                # Send the PDF file
                return send_file(
                    pdf_path,
                    as_attachment=True,
                    download_name=f"{base_filename}_summary.pdf",
                    mimetype='application/pdf'
                )
            except Exception as e:
                logger.error(f"Error creating PDF: {str(e)}")
                return jsonify({'error': f'Failed to create PDF: {str(e)}'}), 500
            finally:
                # Schedule cleanup after response is sent
                if pdf_path:
                    threading.Timer(5.0, lambda: remove_temp_file(pdf_path)).start()
        else:
            # Return the summary as plain text (markdown)
            return summary, 200, {'Content-Type': 'text/plain; charset=utf-8'}
        
    except Exception as e:
        logger.error(f"Unexpected error in upload_file: {str(e)}")
        return jsonify({'error': f'An unexpected error occurred: {str(e)}'}), 500
    finally:
        # Force garbage collection
        gc.collect()

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'message': 'Nexus Notes API is running',
        'timestamp': datetime.now().isoformat(),
        'temp_files_count': len(temp_files)
    })

@app.route('/cleanup', methods=['POST'])
def manual_cleanup():
    """Manual cleanup endpoint"""
    cleanup_temp_files()
    gc.collect()
    return jsonify({'message': 'Cleanup completed', 'temp_files_count': len(temp_files)})

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 5MB.'}), 413

@app.errorhandler(500)
def internal_server_error(e):
    logger.error(f"Internal server error: {str(e)}")
    cleanup_temp_files()  # Clean up on errors
    return jsonify({'error': 'Internal server error. Please try again later.'}), 500

# Background cleanup task
def background_cleanup():
    """Background task to clean up old temp files"""
    while True:
        try:
            time.sleep(300)  # Run every 5 minutes
            cleanup_temp_files()
            gc.collect()
        except Exception as e:
            logger.error(f"Background cleanup error: {str(e)}")

# Start background cleanup thread
cleanup_thread = threading.Thread(target=background_cleanup, daemon=True)
cleanup_thread.start()

if __name__ == '__main__':
    # Check if OpenAI API key is set
    if not os.getenv('OPENAI_API_KEY'):
        print("ERROR: OPENAI_API_KEY environment variable is not set!")
        print("Please set your OpenAI API key as an environment variable:")
        print("export OPENAI_API_KEY='your-api-key-here'")
        exit(1)
    
    print("Starting Nexus Notes API server...")
    print("Make sure to set your OPENAI_API_KEY environment variable")
    app.run(host='0.0.0.0', port=5000, debug=False)  # Disabled debug mode for production